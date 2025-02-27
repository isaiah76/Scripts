from pptx import Presentation
import sys
import os

def extract_text_from_pptx(file_path):
    """Extract text from a PowerPoint presentation (.pptx)."""
    pptx_text = []
    
    try:
        prs = Presentation(file_path)
        for slide_number, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            if slide_text:
                pptx_text.append(f"[Slide {slide_number}]\n" + "\n".join(slide_text))

    except Exception as e:
        raise IOError(f"Error reading PowerPoint file: {e}")

    return "\n\n".join(pptx_text)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python extract_ppt_text.py <file.pptx>")
        sys.exit(1)

    file_path = sys.argv[1]

    if not os.path.exists(file_path):
        print(f"Error: File '{file_path}' not found.")
        sys.exit(1)

    if not file_path.lower().endswith(".pptx"):
        print("Error: Only .pptx files are supported.")
        sys.exit(1)

    extracted_text = extract_text_from_pptx(file_path)
    output_file = os.path.splitext(file_path)[0] + ".txt"
    with open(output_file, "w", encoding="utf-8") as f:
        f.write(extracted_text)

    print(f"Extracted text saved: {output_file}")

